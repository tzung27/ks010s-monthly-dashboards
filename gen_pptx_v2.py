# KS010S 2026-04 Monthly Report Presentation — Premium Modern Design
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── Core helpers ────────────────────────────────────────────────────────────

def rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

i = Inches

def bg(slide, color):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = rgb(color)

def rect(slide, l, t, w, h, fill=None, stroke=None, sw=Pt(0), rnd=False):
    sh = slide.shapes.add_shape(5 if rnd else 1, i(l),i(t),i(w),i(h))
    if rnd:
        try: sh.adjustments[0] = 0.08
        except: pass
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = rgb(fill)
    else: sh.fill.background()
    if stroke: sh.line.width = sw; sh.line.color.rgb = rgb(stroke)
    else: sh.line.fill.background()
    return sh

def oval(slide, l, t, w, h, fill=None, stroke=None, sw=Pt(2)):
    sh = slide.shapes.add_shape(9, i(l),i(t),i(w),i(h))
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = rgb(fill)
    else: sh.fill.background()
    if stroke: sh.line.width = sw; sh.line.color.rgb = rgb(stroke)
    else: sh.line.fill.background()
    return sh

def T(slide, text, l, t, w, h, sz=14, bold=False, clr='#111827',
      align=PP_ALIGN.LEFT, italic=False, fn='微軟正黑體'):
    tb = slide.shapes.add_textbox(i(l),i(t),i(w),i(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = rgb(clr)
    r.font.name = fn
    return tb

def add_line(slide, l, t, w, h, clr, sw=Pt(1.5)):
    from pptx.util import Emu
    ln = slide.shapes.add_connector(1, i(l),i(t), i(l+w),i(t+h))
    ln.line.color.rgb = rgb(clr); ln.line.width = sw
    return ln

# ─── Palette ─────────────────────────────────────────────────────────────────
N  = '#0B1929'   # deep navy (primary dark)
N2 = '#1A3460'   # medium navy
GO = '#B8861C'   # gold (primary accent)
GL = '#F0D585'   # light gold
TL = '#0891B2'   # teal
SK = '#BAE6FD'   # sky blue (light teal)
GR = '#059669'   # green
BL = '#1D4ED8'   # blue
PU = '#7C3AED'   # purple
RD = '#DC2626'   # red
AM = '#D97706'   # amber
WH = '#FFFFFF'
LT = '#F5F8FD'   # very light bg
L2 = '#EAF0FB'   # light blue tint
BD = '#D1DCF0'   # border

MC = {'Aaron':BL,'Denny':GR,'William':PU,'Lucas':AM,'Jimmy':RD}
TEXT2 = '#374151'

# ─── Setup ───────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = i(13.33)
prs.slide_height = i(7.5)
blank = prs.slide_layouts[6]

# ─────────────────────────────────────────────────────────────────────────────
# HELPERS for common patterns
# ─────────────────────────────────────────────────────────────────────────────

def slide_chrome(slide, title, subtitle, accent=TL):
    """Standard header for content slides"""
    # Top gold bar
    rect(slide, 0, 0, 13.33, 0.065, fill=GO)
    # Nav strip
    rect(slide, 0, 0.065, 13.33, 1.18, fill=N)
    # Title
    T(slide, title, 0.42, 0.1, 9.5, 0.72, sz=30, bold=True, clr=WH)
    # Subtitle
    T(slide, subtitle, 0.42, 0.8, 9.5, 0.38, sz=12, clr=SK)
    # Accent line
    rect(slide, 0.42, 1.27, 1.0, 0.055, fill=accent)

def slogan_box(slide, text, y=6.7, bg_col=L2, txt_col=N2, accent_col=GO):
    """Bottom slogan bar"""
    rect(slide, 0, y, 13.33, 0.78, fill=bg_col)
    rect(slide, 0, y, 0.1, 0.78, fill=accent_col)
    T(slide, f'❝  {text}  ❞', 0.28, y+0.12, 12.8, 0.55,
      sz=13, italic=True, clr=N2, align=PP_ALIGN.CENTER)

def stat_card(slide, l, t, w, h, value, unit, label, col, bg_col=WH):
    rect(slide, l, t, w, h, fill=bg_col, stroke=col, sw=Pt(1.5), rnd=True)
    # top color bar inside card
    rect(slide, l, t, w, 0.07, fill=col, rnd=False)
    T(slide, value, l+0.1, t+0.15, w-0.2, h*0.52, sz=34, bold=True, clr=col, align=PP_ALIGN.CENTER)
    T(slide, unit,  l+0.1, t+h*0.52+0.12, w-0.2, 0.28, sz=10, clr='#6B7280', align=PP_ALIGN.CENTER)
    T(slide, label, l+0.1, t+h-0.38, w-0.2, 0.32, sz=10.5, clr=N, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — COVER
# ─────────────────────────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(blank)
bg(s1, LT)

# Left dark panel
rect(s1, 0, 0, 5.45, 7.5, fill=N)
rect(s1, 0, 0, 5.45, 0.1, fill=GO)       # gold top stripe

# Right bg subtle tint
rect(s1, 5.45, 0, 7.88, 7.5, fill=L2)

# Decorative circles (right side, behind content)
oval(s1, 7.2, -0.5, 6.3, 6.3, stroke='#C8D8F0', sw=Pt(1.2))      # large outline circle
oval(s1, 9.8, 4.2, 3.8, 3.8, fill='#DDE8F8')                       # filled bg circle
oval(s1, 5.8, 1.0, 1.1, 1.1, fill='#C9922A')                       # gold dot accent
oval(s1, 7.0, 0.5, 0.45, 0.45, fill='#0891B2')                     # teal dot

# ── Left panel content ──
# KS010S badge
rect(s1, 0.45, 0.25, 1.8, 0.38, fill=GO, rnd=True)
T(s1, 'KS010S  部門', 0.45, 0.25, 1.8, 0.38, sz=11, bold=True, clr=WH, align=PP_ALIGN.CENTER)

# Date
T(s1, '2026  年  4  月', 0.45, 0.8, 4.6, 0.55, sz=19, clr=SK)

# Main title
T(s1, '部門', 0.45, 1.35, 4.7, 0.85, sz=60, bold=True, clr=WH)
T(s1, '工作月報', 0.45, 2.1, 4.7, 0.85, sz=60, bold=True, clr=WH)

# Gold divider
rect(s1, 0.45, 3.08, 3.8, 0.055, fill=GO)

# Slogan lines
T(s1, '以數據為基，以 AI 為翼', 0.45, 3.22, 4.7, 0.45, sz=15.5, italic=True, clr='#A8C4E0')
T(s1, '超越每一個里程碑', 0.45, 3.62, 4.7, 0.45, sz=15.5, italic=True, clr=GL)

# Info line
T(s1, '主管月會 ‧ 總經理簡報', 0.45, 4.2, 4.6, 0.38, sz=11, clr='#5A7A9A')

# AI badge
rect(s1, 0.45, 5.15, 4.55, 0.62, fill='#071E14', stroke='#34D399', sw=Pt(1.8), rnd=True)
oval(s1, 0.6, 5.22, 0.45, 0.45, fill='#34D399')
T(s1, '🤖', 0.6, 5.22, 0.45, 0.45, sz=14, align=PP_ALIGN.CENTER)
T(s1, 'AI 驅動', 1.12, 5.22, 1.5, 0.26, sz=10.5, bold=True, clr='#34D399')
T(s1, '報告 × 儀表板 × 簡報  全程一條龍', 1.12, 5.45, 3.7, 0.28, sz=9.5, clr='#6EE7B7')

# Date bottom
T(s1, '2026.05', 0.45, 7.1, 3, 0.32, sz=11, clr='#334D6B')

# ── Right panel stats (2×2) ──
sw_, sh_ = 3.55, 1.5
stats = [
    ('379',     '件',  '當月服務案件', BL, '#EEF4FF'),
    ('75,000',  'NT$', '實際收費金額', GR, '#ECFDF5'),
    ('216,500', 'NT$', '虛擬服務估值', GO, '#FFFBE8'),
    ('5',       '位',  '部門服務專員', PU, '#F6F3FF'),
]
for ii,(val,unit,lbl,col,bcol) in enumerate(stats):
    row, col_i = ii//2, ii%2
    sx = 5.75 + col_i*(sw_+0.22)
    sy = 1.15 + row*(sh_+0.22)
    stat_card(s1, sx, sy, sw_, sh_, val, unit, lbl, col, bcol)

# Dept label bottom-right
T(s1, 'KS010S · 2026年4月 工作月報', 5.75, 7.1, 7.4, 0.32, sz=10, clr='#6B7280', align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — 執行摘要
# ─────────────────────────────────────────────────────────────────────────────
s2 = prs.slides.add_slide(blank)
bg(s2, WH)
slide_chrome(s2, '執行摘要', '2026年4月 ‧ 部門整體表現一覽')

# Slogan
rect(s2, 0.42, 1.42, 12.5, 0.72, fill='#F5F8FF', stroke=BD, sw=Pt(0.8), rnd=True)
rect(s2, 0.42, 1.42, 0.12, 0.72, fill=GO)
T(s2, '每一件案件，都是客戶信任的具體展現；每一筆數字，都是團隊努力的真實印記',
  0.65, 1.52, 12.1, 0.52, sz=13.5, italic=True, clr=N2, align=PP_ALIGN.CENTER)

# 4 stat boxes
sw2, sh2 = 2.88, 1.35
stats2 = [
    ('379','件','當月服務案件',BL,'#EEF4FF'),
    ('75,000','NT$','實際收費金額',GR,'#ECFDF5'),
    ('216,500','NT$','虛擬服務估值',GO,'#FFFBE8'),
    ('3+','產品線','服務技術範疇',TL,'#EEF9FC'),
]
for ii,(val,unit,lbl,col,bcol) in enumerate(stats2):
    sx = 0.42 + ii*(sw2+0.22)
    stat_card(s2, sx, 2.28, sw2, sh2, val, unit, lbl, col, bcol)

# Highlights title
T(s2, '▌ 本月五大亮點', 0.42, 3.78, 5, 0.38, sz=13, bold=True, clr=N)

hl = [
    ('🏆', WH, GO,  'NT$75,000 實收',
     'William 完成奇美醫院 iOS 高階培訓，創造本月最高實際收費，客戶滿意度極佳'),
    ('📱', WH, GR,  '117 件裝置管理',
     'Denny 單月處理量最高，彰化銀行 / 中華航空大量 iPad DFU 圓滿完成'),
    ('☁️', WH, TL,  'Azure Sentinel POC',
     'Lucas 成功開通 Sentinel 安全監控 POC，為政府與金融客戶資安升級奠定基礎'),
    ('🎨', WH, RD,  'Firefly AI 應用講座',
     'Jimmy 於雲科大舉辦 AI 設計講座並發布教學影片，Adobe 品牌聲量顯著提升'),
    ('🤖', '#071E14', '#34D399', 'AI 月報一條龍',
     '本月報告、互動儀表板、本簡報均由 AI 工具協助完成，展現部門數位轉型實力'),
]
for ii,(icon,bg_c,ic,title,desc) in enumerate(hl):
    ry = 4.22 + ii*0.47
    rect(s2, 0.42, ry, 12.5, 0.44, fill=WH if ii<4 else '#071E14',
         stroke=BD if ii<4 else '#34D399', sw=Pt(0.6), rnd=True)
    # icon bg
    rect(s2, 0.42, ry, 0.44, 0.44, fill=ic, rnd=False)
    T(s2, icon, 0.42, ry, 0.44, 0.44, sz=14, align=PP_ALIGN.CENTER)
    T(s2, title, 0.94, ry+0.05, 2.0, 0.32, sz=11.5, bold=True,
      clr='#34D399' if ii==4 else N)
    T(s2, desc, 3.0, ry+0.05, 9.8, 0.32, sz=10.5, clr='#6EE7B7' if ii==4 else TEXT2)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — 個人績效
# ─────────────────────────────────────────────────────────────────────────────
s3 = prs.slides.add_slide(blank)
bg(s3, LT)
slide_chrome(s3, '個人績效一覽', '5 位服務專員 ‧ 案件數 ‧ 虛擬費用 ‧ 實收費用')

members3 = [
    ('Aaron',  '全能型\n雲端支援', 96, 'NT$51,100', 'NT$—',   BL, '#EEF4FF'),
    ('Denny',  '端點管理\n資安主力',117,'NT$83,200', 'NT$—',   GR, '#ECFDF5'),
    ('William','高價值\n服務貢獻', 43, 'NT$52,800', 'NT$75,000',PU,'#F6F3FF'),
    ('Lucas',  'Azure\n雲端架構', 51, 'NT$29,400', 'NT$—',   AM, '#FFFBE8'),
    ('Jimmy',  'Adobe/Corel\n負責人', 72, '—',    '—',        RD, '#FFF2F2'),
]

cw = 2.44
for ii,(name,role,cases,virt,act,col,bcol) in enumerate(members3):
    cx = 0.28 + ii*(cw+0.18)
    cy = 1.5

    # Card bg
    rect(s3, cx, cy, cw, 5.55, fill=bcol, stroke=col, sw=Pt(1.5), rnd=True)

    # Colored header
    rect(s3, cx, cy, cw, 1.0, fill=col, rnd=False)
    # Name
    T(s3, name, cx+0.12, cy+0.06, cw-0.24, 0.48,
      sz=20, bold=True, clr=WH, align=PP_ALIGN.LEFT)
    # Role (multiline via \n)
    role_lines = role.replace('\n',' · ')
    T(s3, role_lines, cx+0.1, cy+0.54, cw-0.2, 0.42, sz=9, clr=WH)

    # Cases: big number
    T(s3, str(cases), cx+0.1, cy+1.08, cw-0.2, 0.72,
      sz=44, bold=True, clr=col, align=PP_ALIGN.CENTER)
    T(s3, '件', cx+0.1, cy+1.76, cw-0.2, 0.28,
      sz=11, clr='#6B7280', align=PP_ALIGN.CENTER)

    # Divider
    rect(s3, cx+0.2, cy+2.12, cw-0.4, 0.04, fill=BD)

    # Virtual fee
    rect(s3, cx+0.18, cy+2.28, cw-0.36, 0.88, fill=WH, stroke=BD, sw=Pt(0.8), rnd=True)
    T(s3, '虛擬費用', cx+0.28, cy+2.34, cw-0.56, 0.25, sz=8.5, clr='#9CA3AF')
    T(s3, virt, cx+0.22, cy+2.57, cw-0.44, 0.48,
      sz=13.5 if len(virt)>5 else 15, bold=True, clr=col)

    # Actual fee
    act_col = '#059669' if 'NT$' in act and '—' not in act else '#9CA3AF'
    rect(s3, cx+0.18, cy+3.26, cw-0.36, 0.88, fill=WH, stroke=BD, sw=Pt(0.8), rnd=True)
    T(s3, '實際收費', cx+0.28, cy+3.32, cw-0.56, 0.25, sz=8.5, clr='#9CA3AF')
    T(s3, act, cx+0.22, cy+3.55, cw-0.44, 0.48,
      sz=13.5 if len(act)>5 else 15, bold=True, clr=act_col)

    # William star badge
    if name == 'William':
        rect(s3, cx+0.18, cy+4.28, cw-0.36, 0.5,
             fill='#FEF3C7', stroke='#F59E0B', sw=Pt(1.2), rnd=True)
        T(s3, '★  本月最高實收', cx+0.18, cy+4.3, cw-0.36, 0.44,
          sz=10, bold=True, clr='#92400E', align=PP_ALIGN.CENTER)

slogan_box(s3, '個人卓越，成就團隊；團隊卓越，創造客戶無可取代的價值')

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — 本月重大里程碑
# ─────────────────────────────────────────────────────────────────────────────
s4 = prs.slides.add_slide(blank)
bg(s4, WH)
slide_chrome(s4, '本月重大里程碑', '代表性成就 ‧ 高價值案件 ‧ 客戶突破')

milestones = [
    ('Aaron', BL, '96件',
     '雲端遷移\n顧問服務',
     ['Cloudiway 跨租戶\n郵件遷移諮詢',
      'GWS → M365\n多家客戶遷移輔導',
      'Entra ID 取代 AD\n架構規劃完成']),
    ('Denny', GR, '117件',
     '裝置管理\n最高產出',
     ['彰化銀行\n大量 iPad DFU ✓',
      '中華航空\nADE 裝置納管 ✓',
      'Trend Micro\nApex One 建置完成']),
    ('William', PU, '43件',
     'NT$75,000\n實收最高',
     ['奇美醫院\niOS 高階培訓完成',
      'Kaspersky 17件\n資安管理案件',
      '多家客戶\n擴展合約準備中']),
    ('Lucas', AM, '51件',
     'Azure 安全\n架構升級',
     ['Sentinel POC\n正式開通 ✓',
      'GPU 配額\n提升申請成功',
      '坤侑科技\nAWS→Azure 遷移啟動']),
    ('Jimmy', RD, '72件',
     'AI 教育\n品牌拓展',
     ['雲科大\nFirefly AI 講座 ✓',
      'AI 向量設計\n教學影片上線',
      'Acrobat Sign\n企業導入支援']),
]

cw2 = 2.44
for ii,(name,col,badge,sub,pts) in enumerate(milestones):
    cx = 0.28 + ii*(cw2+0.18)
    cy = 1.55

    # Badge
    rect(s4, cx, cy, cw2, 0.38, fill=col, rnd=True)
    T(s4, f'{name}  {badge}', cx+0.1, cy+0.04, cw2-0.2, 0.3,
      sz=12.5, bold=True, clr=WH, align=PP_ALIGN.CENTER)

    # Sub-role
    rect(s4, cx, cy+0.38, cw2, 0.6, fill='#F8FAFC', stroke=col, sw=Pt(1.2))
    T(s4, sub.replace('\n',' '), cx+0.1, cy+0.42, cw2-0.2, 0.52,
      sz=11, bold=True, clr=col, align=PP_ALIGN.CENTER)

    # Milestones body
    rect(s4, cx, cy+0.98, cw2, 4.62, fill=WH, stroke=BD, sw=Pt(1))
    for jj, pt in enumerate(pts):
        py = cy+0.98 + jj*1.44
        # Number circle
        oval(s4, cx+0.2, py+0.2, 0.45, 0.45, fill=col)
        T(s4, str(jj+1), cx+0.2, py+0.2, 0.45, 0.45,
          sz=11, bold=True, clr=WH, align=PP_ALIGN.CENTER)
        T(s4, pt.replace('\n',' '), cx+0.72, py+0.15, cw2-0.82, 0.68,
          sz=11, clr=N)

slogan_box(s4, '突破不是偶然，是每日專業積累、每次客戶信任的印記')

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — 5月工作重點
# ─────────────────────────────────────────────────────────────────────────────
s5 = prs.slides.add_slide(blank)
bg(s5, LT)
slide_chrome(s5, '5 月工作重點預覽', '各成員下月重點任務 ‧ 持續衝刺')

next5 = [
    ('Aaron', BL, [
        'Cloudiway T2T\n跨租戶遷移實作',
        'Defender for Business\n部署案件跟進',
        'Entra ID 取代 AD\n客戶正式規劃',
        'Power Automate\nSharePoint 整合開發',
    ]),
    ('Denny', GR, [
        '彰銀 Jamf\n憑證更新前置作業',
        '強茂 Jamf Mobile\nPOC 後續跟進',
        '凱基人壽 macOS\nMDM 服務擴展',
        'Apex One 三鶯線\n建置完成驗收',
    ]),
    ('William', PU, [
        '華航 Jamf\n年度服務續約洽談',
        '奇美醫院\n進階培訓規劃',
        'Kaspersky KES\n升版服務推廣',
        '強茂 MDM POC\n轉正式合約',
    ]),
    ('Lucas', AM, [
        'Sentinel POC\n驗收報告撰寫',
        'Azure 備份 SOP\n標準化建立',
        '坤侑科技\nAWS→Azure 遷移規劃',
        '嘉龍資訊 VM\n升版完成驗收',
    ]),
    ('Jimmy', RD, [
        'Firefly AI 進階\n課程規劃上線',
        'Acrobat Sign\n企業版推廣',
        'AI 工具整合\n培訓教材製作',
        '技術文件月刊\n部落格定期產出',
    ]),
]

cw3 = 2.44
for ii,(name,col,tasks) in enumerate(next5):
    cx = 0.28 + ii*(cw3+0.18)
    cy = 1.55
    # Header
    rect(s5, cx, cy, cw3, 0.52, fill=col, rnd=True)
    T(s5, name, cx+0.1, cy+0.07, cw3-0.2, 0.38,
      sz=17, bold=True, clr=WH, align=PP_ALIGN.CENTER)
    # Body
    rect(s5, cx, cy+0.52, cw3, 5.08, fill=WH, stroke=col, sw=Pt(1.3))
    for jj,task in enumerate(tasks):
        ty = cy+0.52 + jj*1.22
        # Number tag
        rect(s5, cx+0.15, ty+0.2, 0.35, 0.35, fill=col, rnd=True)
        T(s5, str(jj+1), cx+0.15, ty+0.2, 0.35, 0.35,
          sz=10, bold=True, clr=WH, align=PP_ALIGN.CENTER)
        T(s5, task.replace('\n',' '), cx+0.56, ty+0.12, cw3-0.68, 0.68,
          sz=11.5, clr=N)
        if jj < 3:
            rect(s5, cx+0.15, ty+1.12, cw3-0.3, 0.04, fill='#E8F0F8')

slogan_box(s5, '今日佈局，定義明日格局；持續執行，方能引領市場')

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — AI 月報一條龍 (SHOWPIECE)
# ─────────────────────────────────────────────────────────────────────────────
s6 = prs.slides.add_slide(blank)
bg(s6, '#070F1F')

# Decorative bg elements
oval(s6, -1, -1, 5.5, 5.5, stroke='#1A3A5C', sw=Pt(2))
oval(s6, 9.8, 3.5, 5.5, 5.5, stroke='#0A2840', sw=Pt(1.5))
oval(s6, 10.5, -0.5, 3.5, 3.5, fill='#0A1E2E')

# Top accent bar (green-themed)
rect(s6, 0, 0, 13.33, 0.065, fill='#059669')

# Title area
rect(s6, 0, 0.065, 13.33, 1.28, fill='#071E14')
T(s6, 'AI 驅動月報一條龍', 0.42, 0.1, 8, 0.72, sz=30, bold=True, clr='#34D399')
T(s6, '從原始資料到管理簡報，全程 AI 工具協助完成，零人工排版，數據驅動決策',
  0.42, 0.8, 11, 0.42, sz=12, clr='#86EFAC')
rect(s6, 0.42, 1.32, 1.0, 0.055, fill='#34D399')

# Core quote
rect(s6, 0.42, 1.52, 12.5, 0.68, fill='#0A2010', stroke='#34D399', sw=Pt(1.2), rnd=True)
T(s6, '❝  AI 不是工具，而是新一代的工作夥伴；讓人從重複勞動中解放，聚焦真正的創造力  ❞',
  0.62, 1.6, 12.1, 0.52, sz=13, italic=True, clr='#86EFAC', align=PP_ALIGN.CENTER)

# Flow steps
flow = [
    ('📊', '原始資料', 'Excel 工作紀錄\nPDF 月報文件', '#0C2A4A', '#3B82F6'),
    ('🤖', 'AI 分析', 'Claude AI 提取\n結構化資料', '#064E3B', '#34D399'),
    ('📈', '互動儀表板', '自動生成\nWeb Dashboard', '#2D1460', '#A855F7'),
    ('📋', '部門月報', 'KPI 追蹤\n案件明細一覽', '#451A03', '#F59E0B'),
    ('🎯', '本份簡報', 'AI 自動排版\n簡報一條龍', '#450A0A', '#EF4444'),
]
fw = 2.3
for ii,(icon,title,desc,bg_c,col) in enumerate(flow):
    fx = 0.42 + ii*(fw+0.2)
    fy = 2.38
    rect(s6, fx, fy, fw, 2.32, fill=bg_c, stroke=col, sw=Pt(1.5), rnd=True)
    # icon circle
    oval(s6, fx+fw/2-0.3, fy+0.12, 0.6, 0.6, fill=col)
    T(s6, icon, fx+fw/2-0.3, fy+0.12, 0.6, 0.6, sz=16, align=PP_ALIGN.CENTER)
    T(s6, title, fx+0.1, fy+0.82, fw-0.2, 0.42,
      sz=14.5, bold=True, clr=col, align=PP_ALIGN.CENTER)
    T(s6, desc, fx+0.12, fy+1.26, fw-0.24, 0.78,
      sz=10.5, clr='#94A3B8', align=PP_ALIGN.CENTER)
    # Arrow connector
    if ii < 4:
        T(s6, '▶', fx+fw+0.04, fy+0.95, 0.18, 0.38,
          sz=14, clr='#334155', align=PP_ALIGN.CENTER)

# Benefits row
benefits = [
    ('⏱️', '效率提升 80%', '報告製作時間大幅縮短'),
    ('🎨', '零人工排版', '自動生成，無需手動美化'),
    ('🔄', '即時更新', '資料異動後一鍵重新生成'),
    ('📐', '數據驅動', '視覺化圖表自動彙整分析'),
]
bw6 = 2.88
for ii,(icon,title,desc) in enumerate(benefits):
    bx = 0.6 + ii*(bw6+0.34)
    rect(s6, bx, 5.0, bw6, 0.98, fill='#0F1F38', stroke='#1E3A5C', sw=Pt(1), rnd=True)
    T(s6, icon, bx+0.12, 5.06, 0.42, 0.42, sz=18)
    T(s6, title, bx+0.58, 5.08, bw6-0.7, 0.38, sz=12, bold=True, clr='#34D399')
    T(s6, desc, bx+0.58, 5.44, bw6-0.7, 0.42, sz=10, clr='#64748B')

# Footer
rect(s6, 0, 6.12, 13.33, 1.38, fill='#040D18')
T(s6, '🤖  本月報告 ‧ 互動儀表板 ‧ 本份簡報  均由 AI 工具（Claude）協助完成',
  0.42, 6.24, 12.5, 0.42, sz=12, bold=True, clr='#34D399', align=PP_ALIGN.CENTER)
T(s6, 'KS010S 部門 AI-First 工作模式實踐  ‧  2026年4月',
  0.42, 6.65, 12.5, 0.35, sz=10.5, clr='#334D6B', align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — 年度 KPI 進度
# ─────────────────────────────────────────────────────────────────────────────
s7 = prs.slides.add_slide(blank)
bg(s7, WH)
slide_chrome(s7, '2026 年度 KPI 進度', '目標追蹤 ‧ 進度健康 ‧ 持續聚焦')

# Slogan
rect(s7, 0.42, 1.42, 12.5, 0.58, fill='#FFFBEB', stroke='#FCD34D', sw=Pt(0.8), rnd=True)
rect(s7, 0.42, 1.42, 0.1, 0.58, fill=GO)
T(s7, '目標是方向，執行是力量，結果是最有力的証明',
  0.65, 1.5, 12.1, 0.42, sz=13, italic=True, clr=N2, align=PP_ALIGN.CENTER)

kpis = [
    ('📞', '月均技術支援案件達標', GR, '已達成',
     'Denny 117件 ✓ · William 43件 ✓ · Lucas 51件 ✓ · Jimmy 72件 ✓'),
    ('💰', '收費服務成交目標', BL, '進行中',
     'William NT$75,000 實收 ✓ · 其他成員積極推進收費案件中'),
    ('🎤', '對外 Workshop 年度目標 ≥8次', AM, '進行中',
     'Jimmy 雲科大 Firefly AI 講座已完成 ✓ · 年度累計持續推進中'),
    ('📚', '知識庫文件撰寫', TL, '進行中',
     'Jimmy KB 文章定期更新 ✓ · Denny / William 季度目標推進中'),
    ('🏆', '原廠認證更新', PU, '待執行',
     '各成員按代理商要求維持有效認證，年度前完成 renew'),
    ('🤖', 'AI 工具導入', GR, '已達成',
     'Jimmy Firefly AI 整合推廣 ✓ · William / Denny 案件回覆自動化 ✓'),
    ('📊', '部門月報按時呈報', GR, '已達成',
     '每月 AI 協助彙整，互動儀表板 + 簡報雙軌呈報 ✓'),
]

status_cfg = {
    '已達成': ('#D1FAE5','#065F46','#ECFDF5'),
    '進行中': ('#DBEAFE','#1D4ED8','#EFF6FF'),
    '待執行': ('#FEF3C7','#92400E','#FFFBEB'),
}
for ii,(icon,title,col,status,detail) in enumerate(kpis):
    ry = 2.15 + ii*0.65
    sb, sc, rb = status_cfg.get(status, ('#F1F5F9','#475569','#F8FAFC'))
    rect(s7, 0.42, ry, 12.5, 0.58, fill=rb, stroke=BD, sw=Pt(0.5), rnd=True)
    # Left icon bg
    rect(s7, 0.42, ry, 0.58, 0.58, fill=col, rnd=False)
    T(s7, icon, 0.42, ry, 0.58, 0.58, sz=15, align=PP_ALIGN.CENTER, clr=WH)
    # Title
    T(s7, title, 1.1, ry+0.08, 3.5, 0.38, sz=12, bold=True, clr=N)
    # Detail
    T(s7, detail, 4.75, ry+0.1, 7.4, 0.38, sz=10.5, clr=TEXT2)
    # Status badge
    rect(s7, 12.2, ry+0.1, 0.65, 0.35, fill=sb, rnd=True)
    T(s7, status, 12.2, ry+0.1, 0.65, 0.35, sz=9, bold=True, clr=sc, align=PP_ALIGN.CENTER)

slogan_box(s7, '每一項 KPI，都是對客戶、對公司、對自己的承諾', bg_col='#F5F8FF')

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — 結語
# ─────────────────────────────────────────────────────────────────────────────
s8 = prs.slides.add_slide(blank)
bg(s8, N)

# Decorative circles
oval(s8, -2, -1, 7, 7, stroke='#1A3460', sw=Pt(1.5))
oval(s8, 8.5, 3, 6, 6, fill='#0C1E38')
oval(s8, 11.5, -0.5, 4, 4, stroke='#B8861C', sw=Pt(1.5))

# Gold top stripe
rect(s8, 0, 0, 13.33, 0.1, fill=GO)

# Large title
T(s8, '持續進化', 0.6, 0.5, 8, 1.1, sz=62, bold=True, clr=WH)
T(s8, '方能引領未來', 0.6, 1.5, 8, 0.9, sz=50, bold=True, clr=GL)

# Divider
rect(s8, 0.6, 2.52, 5.5, 0.07, fill=GO)

# Sub-title
T(s8, 'KS010S 部門  ·  2026年4月 月報圓滿完成', 0.6, 2.72, 8, 0.42, sz=14, clr=SK)

# 4 key takeaways
takes = [
    (BL, '業績健康', '379件服務，NT$75K 實收，虛擬估值 NT$216K'),
    (GR, 'AI 落地', '月報、儀表板、簡報一條龍成功實踐'),
    (GO, '佈局 5月', '彰銀、奇美、Sentinel、Firefly AI 等重點項目持續推進'),
    (PU, 'KPI 達標', '月均案件、AI 導入、月報呈報三項已達標'),
]
for ii,(col,title,desc) in enumerate(takes):
    tx = 0.6 + ii*3.1
    rect(s8, tx, 3.5, 2.85, 1.55, fill='#0D2345', stroke=col, sw=Pt(1.5), rnd=True)
    rect(s8, tx, 3.5, 2.85, 0.07, fill=col)
    T(s8, title, tx+0.15, 3.62, 2.55, 0.45, sz=15, bold=True, clr=col)
    T(s8, desc, tx+0.15, 4.04, 2.55, 0.78, sz=10.5, clr='#94A3B8')

# Thank you
T(s8, 'Thank you', 0.6, 5.4, 8, 0.72, sz=36, bold=True, clr='#1E3A6E', italic=True, fn='Times New Roman')

# AI badge bottom-right
rect(s8, 8.5, 5.25, 4.55, 0.68, fill='#071E14', stroke='#34D399', sw=Pt(1.5), rnd=True)
T(s8, '🤖  本月報告 · 儀表板 · 簡報', 8.65, 5.3, 4.25, 0.28, sz=11, bold=True, clr='#34D399')
T(s8, '全程由 AI 工具（Claude）協助完成', 8.65, 5.58, 4.25, 0.28, sz=10, clr='#6EE7B7')

# Footer line
rect(s8, 0, 7.15, 13.33, 0.35, fill='#060E1C')
T(s8, 'KS010S 部門月報  ‧  2026年4月  ‧  主管月會  ‧  2026.05',
  0.42, 7.17, 12.5, 0.3, sz=10, clr='#334D6B', align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────────────────────────
out = r'D:\excel_dashboard07 月報\KS010S_2026年4月月報簡報_v2.pptx'
prs.save(out)
print(f'✓ Saved: {out}')
