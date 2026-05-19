from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Helpers ──────────────────────────────────────────────────────────────────

def rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = rgb(line)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color='#1E293B',
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    return tb

def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)

def slide_header(slide, title, subtitle='', bg='#0F172A', accent='#38BDF8'):
    W = Inches(13.33)
    # full-width top bar
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.35), fill=bg)
    add_text(slide, title, Inches(0.45), Inches(0.12), Inches(9), Inches(0.72),
             size=26, bold=True, color='#FFFFFF')
    if subtitle:
        add_text(slide, subtitle, Inches(0.45), Inches(0.82), Inches(9), Inches(0.4),
                 size=13, color=accent)
    # accent line
    add_rect(slide, Inches(0), Inches(1.35), W, Inches(0.04), fill=accent)

def stat_box(slide, l, t, w, h, value, label, val_color='#0F172A', bg='#F8FAFC', border='#E2E8F0'):
    add_rect(slide, l, t, w, h, fill=bg, line=border, line_w=Pt(1.2))
    add_text(slide, value, l+Inches(0.15), t+Inches(0.12), w-Inches(0.3), Inches(0.6),
             size=30, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    add_text(slide, label, l+Inches(0.1), t+Inches(0.68), w-Inches(0.2), Inches(0.35),
             size=11, color='#64748B', align=PP_ALIGN.CENTER)

# ── Presentation Setup ───────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # completely blank

C = {
    'navy':   '#0F172A',
    'teal':   '#0891B2',
    'sky':    '#38BDF8',
    'green':  '#22C55E',
    'amber':  '#F59E0B',
    'red':    '#EF4444',
    'purple': '#A855F7',
    'blue':   '#3B82F6',
    'white':  '#FFFFFF',
    'light':  '#F1F5F9',
    'subtle': '#64748B',
    'dark':   '#1E293B',
}

W = Inches(13.33)
H = Inches(7.5)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Cover
# ─────────────────────────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(blank_layout)
set_slide_bg(s1, '#0F172A')

# gradient strip left
add_rect(s1, Inches(0), Inches(0), Inches(0.18), H, fill='#0891B2')

# department badge top-right
add_rect(s1, Inches(10.6), Inches(0.3), Inches(2.5), Inches(0.55), fill='#0891B2')
add_text(s1, 'KS010S 部門月報', Inches(10.6), Inches(0.3), Inches(2.5), Inches(0.55),
         size=12, bold=True, color='#FFFFFF', align=PP_ALIGN.CENTER)

# main title
add_text(s1, '2026 年 4 月', Inches(0.55), Inches(1.4), Inches(10), Inches(0.7),
         size=22, color='#38BDF8')
add_text(s1, '部門工作月報', Inches(0.55), Inches(2.0), Inches(10), Inches(1.3),
         size=58, bold=True, color='#FFFFFF')
add_text(s1, '主管月會 ‧ 總經理簡報', Inches(0.55), Inches(3.28), Inches(10), Inches(0.55),
         size=20, color='#94A3B8')

# divider
add_rect(s1, Inches(0.55), Inches(3.9), Inches(5.5), Inches(0.04), fill='#38BDF8')

# key numbers row
stats = [('379件', '當月服務案件'), ('NT$75,000', '實收費用'), ('NT$216,500', '虛擬服務估值'), ('5位', '服務專員')]
bw = Inches(2.6)
for i,(v,l) in enumerate(stats):
    bx = Inches(0.55) + i * (bw + Inches(0.22))
    add_rect(s1, bx, Inches(4.1), bw, Inches(0.95), fill='#1E293B', line='#334155', line_w=Pt(1))
    add_text(s1, v, bx+Inches(0.12), bx-Inches(3.5), bw-Inches(0.24), Inches(0.52),
             size=20, bold=True, color='#38BDF8', align=PP_ALIGN.CENTER)
    add_text(s1, l, bx+Inches(0.1), Inches(4.65), bw-Inches(0.2), Inches(0.35),
             size=10, color='#94A3B8', align=PP_ALIGN.CENTER)

for i,(v,l) in enumerate(stats):
    bx = Inches(0.55) + i * (bw + Inches(0.22))
    add_text(s1, v, bx+Inches(0.1), Inches(4.15), bw-Inches(0.2), Inches(0.5),
             size=20, bold=True, color='#38BDF8', align=PP_ALIGN.CENTER)
    add_text(s1, l, bx+Inches(0.1), Inches(4.63), bw-Inches(0.2), Inches(0.32),
             size=10, color='#94A3B8', align=PP_ALIGN.CENTER)

# AI badge
add_rect(s1, Inches(0.55), Inches(5.25), Inches(4.5), Inches(0.72), fill='#064E3B', line='#22C55E', line_w=Pt(1.5))
add_text(s1, '🤖  AI 驅動 ‧ 報告 × 儀表板 × 簡報  全程一條龍', Inches(0.7), Inches(5.3), Inches(4.2), Inches(0.6),
         size=12, bold=True, color='#4ADE80', align=PP_ALIGN.LEFT)

# date
add_text(s1, '2026年5月', Inches(0.55), Inches(6.85), Inches(6), Inches(0.4),
         size=12, color='#475569')

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — 執行摘要 / 亮點
# ─────────────────────────────────────────────────────────────────────────────
s2 = prs.slides.add_slide(blank_layout)
set_slide_bg(s2, '#F8FAFC')
slide_header(s2, '執行摘要', '2026 年 4 月 ‧ KS010S 部門亮點速覽')

# 4 stat boxes
boxes = [
    ('379', '當月服務案件總數', '#0891B2', '#EFF6FF'),
    ('NT$75,000', '實際收費金額', '#16A34A', '#F0FDF4'),
    ('NT$216,500', '虛擬服務估值', '#D97706', '#FFFBEB'),
    ('5 人', '部門服務專員', '#7C3AED', '#F5F3FF'),
]
bw = Inches(2.95)
for i,(v,l,vc,bc) in enumerate(boxes):
    bx = Inches(0.3) + i*(bw+Inches(0.2))
    add_rect(s2, bx, Inches(1.55), bw, Inches(1.1), fill=bc, line='#E2E8F0', line_w=Pt(1))
    add_text(s2, v, bx+Inches(0.1), Inches(1.6), bw-Inches(0.2), Inches(0.6),
             size=28, bold=True, color=vc, align=PP_ALIGN.CENTER)
    add_text(s2, l, bx+Inches(0.1), Inches(2.17), bw-Inches(0.2), Inches(0.4),
             size=11, color='#475569', align=PP_ALIGN.CENTER)

# highlights title
add_text(s2, '▌ 本月重大亮點', Inches(0.3), Inches(2.9), Inches(6), Inches(0.45),
         size=14, bold=True, color=C['navy'])

highlights = [
    '🏆  William 完成奇美醫院 iOS 高階培訓，創造本月最高實收費用 NT$75,000',
    '📱  Denny 處理 117 件裝置管理案件，單月最高，彰化銀行 / 中華航空大量 iPad DFU 圓滿完成',
    '☁️  Lucas 成功開通 Sentinel 安全監控 POC，為政府與金融客戶資安升級奠基',
    '🎨  Jimmy 於雲科大舉辦 Firefly AI 應用講座，AI 設計教學影片正式上線 YouTube',
    '🤖  本月報表、互動儀表板、本簡報全程由 AI 工具協助完成，實現月報一條龍',
]
for i, h in enumerate(highlights):
    add_text(s2, h, Inches(0.35), Inches(3.35)+i*Inches(0.63), Inches(12.6), Inches(0.58),
             size=13, color=C['dark'])

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — 個人績效一覽
# ─────────────────────────────────────────────────────────────────────────────
s3 = prs.slides.add_slide(blank_layout)
set_slide_bg(s3, '#F8FAFC')
slide_header(s3, '個人績效一覽', '5位服務專員 ‧ 案件數 ‧ 虛擬費用 ‧ 實收費用')

members = [
    ('Aaron',   '全能型雲端支援專家',  96, 'NT$51,100', 'NT$0',      '#3B82F6', '#EFF6FF'),
    ('Denny',   '端點管理與資安主力', 117, 'NT$83,200', 'NT$0',      '#22C55E', '#F0FDF4'),
    ('William', '高價值服務與營收貢獻者', 43, 'NT$52,800', 'NT$75,000', '#A855F7', '#FAF5FF'),
    ('Lucas',   'Azure 雲端基礎架構手', 51, 'NT$29,400', 'NT$0',      '#F59E0B', '#FFFBEB'),
    ('Jimmy',   'Adobe / Corel 負責人', 72, '—',         '—',         '#EF4444', '#FFF1F2'),
]

card_w = Inches(2.45)
card_h = Inches(4.8)
for i,(name,role,cases,virt,act,col,bg) in enumerate(members):
    cx = Inches(0.25) + i*(card_w+Inches(0.17))
    cy = Inches(1.6)
    add_rect(s3, cx, cy, card_w, card_h, fill=bg, line='#E2E8F0', line_w=Pt(1))
    # colored header
    add_rect(s3, cx, cy, card_w, Inches(0.78), fill=col)
    add_text(s3, name, cx+Inches(0.12), cy+Inches(0.04), card_w-Inches(0.24), Inches(0.42),
             size=18, bold=True, color='#FFFFFF')
    add_text(s3, role, cx+Inches(0.1), cy+Inches(0.44), card_w-Inches(0.2), Inches(0.32),
             size=8.5, color='#FFFFFF')
    # stats
    add_text(s3, str(cases), cx+Inches(0.1), cy+Inches(0.95), card_w-Inches(0.2), Inches(0.58),
             size=34, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s3, '件', cx+Inches(0.1), cy+Inches(1.5), card_w-Inches(0.2), Inches(0.3),
             size=11, color='#64748B', align=PP_ALIGN.CENTER)
    # virtual
    add_rect(s3, cx+Inches(0.15), cy+Inches(1.95), card_w-Inches(0.3), Inches(0.72),
             fill='#FFFFFF', line='#E2E8F0', line_w=Pt(0.8))
    add_text(s3, '虛擬費用', cx+Inches(0.2), cy+Inches(1.98), card_w-Inches(0.4), Inches(0.26),
             size=9, color='#94A3B8')
    add_text(s3, virt, cx+Inches(0.18), cy+Inches(2.22), card_w-Inches(0.36), Inches(0.38),
             size=13, bold=True, color=col)
    # actual
    act_col = '#16A34A' if act != 'NT$0' and act != '—' else '#94A3B8'
    add_rect(s3, cx+Inches(0.15), cy+Inches(2.78), card_w-Inches(0.3), Inches(0.72),
             fill='#FFFFFF', line='#E2E8F0', line_w=Pt(0.8))
    add_text(s3, '實際收費', cx+Inches(0.2), cy+Inches(2.81), card_w-Inches(0.4), Inches(0.26),
             size=9, color='#94A3B8')
    add_text(s3, act, cx+Inches(0.18), cy+Inches(3.05), card_w-Inches(0.36), Inches(0.38),
             size=13, bold=True, color=act_col)

# William star
add_rect(s3, Inches(2.49*2+0.25*3+0.17*2), Inches(1.55), Inches(0.65), Inches(0.35), fill='#FEF3C7', line='#F59E0B', line_w=Pt(1))
add_text(s3, '★ 最高收費', Inches(2.49*2+0.25*3+0.17*2), Inches(1.55), Inches(0.65), Inches(0.35),
         size=7.5, bold=True, color='#92400E', align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — 本月重大里程碑
# ─────────────────────────────────────────────────────────────────────────────
s4 = prs.slides.add_slide(blank_layout)
set_slide_bg(s4, '#F8FAFC')
slide_header(s4, '本月重大里程碑', '各成員代表性成就')

milestones = [
    ('#3B82F6', 'Aaron',   '96件', 'M365 / Exchange / Entra',
     ['Cloudiway 郵件跨租戶遷移諮詢完成', 'GWS → M365 多家遷移輔導', 'Entra ID 取代地端 AD 規劃']),
    ('#22C55E', 'Denny',   '117件', 'Apple & Jamf / Trend Micro',
     ['彰化銀行 大量 iPad DFU 完成', '中華航空 ADE 裝置大量納管', 'Trend Micro Apex One 建置完成']),
    ('#A855F7', 'William', '43件', 'Kaspersky / Jamf Pro / TeamViewer',
     ['奇美醫院 iOS 高階培訓（NT$75K）', 'Kaspersky 17件安全管理', '客戶擴展合約準備中']),
    ('#F59E0B', 'Lucas',   '51件', 'Azure 雲端',
     ['Azure Sentinel POC 正式開通', 'GPU 配額提升申請成功', '坤侑科技 AWS→Azure 遷移規劃啟動']),
    ('#EF4444', 'Jimmy',   '72件', 'Adobe CC / Corel / AI 教育',
     ['雲科大 Firefly AI 應用講座', 'AI 向量設計教學影片上線', 'Acrobat Sign 企業導入']),
]

col_w = Inches(2.5)
for i,(col,name,cnt,prod,pts) in enumerate(milestones):
    cx = Inches(0.2) + i*(col_w+Inches(0.18))
    cy = Inches(1.55)
    # header strip
    add_rect(s4, cx, cy, col_w, Inches(0.6), fill=col)
    add_text(s4, f'{name}  {cnt}', cx+Inches(0.12), cy+Inches(0.06), col_w-Inches(0.24), Inches(0.32),
             size=14, bold=True, color='#FFFFFF')
    add_text(s4, prod, cx+Inches(0.1), cy+Inches(0.36), col_w-Inches(0.2), Inches(0.22),
             size=8.5, color='#FFFFFF')
    # body
    add_rect(s4, cx, cy+Inches(0.6), col_w, Inches(5.0), fill='#FFFFFF', line='#E2E8F0', line_w=Pt(1))
    for j,pt in enumerate(pts):
        add_text(s4, f'• {pt}', cx+Inches(0.15), cy+Inches(0.75)+j*Inches(0.7),
                 col_w-Inches(0.25), Inches(0.65), size=11, color=C['dark'])

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — 下月工作重點
# ─────────────────────────────────────────────────────────────────────────────
s5 = prs.slides.add_slide(blank_layout)
set_slide_bg(s5, '#F8FAFC')
slide_header(s5, '5 月工作重點預覽', '各成員下月重點任務')

next_month = [
    ('Aaron',   '#3B82F6', ['Cloudiway T2T 遷移實作', 'Defender for Business 部署', 'Entra ID 客戶規劃', 'Power Automate 整合開發']),
    ('Denny',   '#22C55E', ['彰銀 Jamf 憑證更新', '強茂 Jamf Mobile POC', '凱基人壽 macOS MDM 擴展', 'Apex One 三鶯線驗收']),
    ('William', '#A855F7', ['華航 Jamf 年約洽談', '奇美醫院進階培訓規劃', 'Kaspersky KES 升版推廣', '強茂 MDM 轉正式合約']),
    ('Lucas',   '#F59E0B', ['Sentinel POC 驗收報告', 'Azure 備份 SOP 標準化', 'AWS→Azure 遷移規劃', '嘉龍資訊 VM 升版驗收']),
    ('Jimmy',   '#EF4444', ['Firefly AI 進階課程規劃', 'Acrobat Sign 企業版推廣', 'AI 工具培訓教材製作', '技術文件月刊產出']),
]

col_w = Inches(2.5)
for i,(name,col,tasks) in enumerate(next_month):
    cx = Inches(0.2) + i*(col_w+Inches(0.18))
    cy = Inches(1.55)
    add_rect(s5, cx, cy, col_w, Inches(0.5), fill=col)
    add_text(s5, name, cx+Inches(0.12), cy+Inches(0.06), col_w-Inches(0.24), Inches(0.38),
             size=15, bold=True, color='#FFFFFF', align=PP_ALIGN.CENTER)
    add_rect(s5, cx, cy+Inches(0.5), col_w, Inches(5.1), fill='#FFFFFF', line='#E2E8F0', line_w=Pt(1))
    for j,task in enumerate(tasks):
        # number badge
        add_rect(s5, cx+Inches(0.15), cy+Inches(0.68)+j*Inches(1.1), Inches(0.32), Inches(0.32),
                 fill=col)
        add_text(s5, str(j+1), cx+Inches(0.15), cy+Inches(0.68)+j*Inches(1.1), Inches(0.32), Inches(0.32),
                 size=9, bold=True, color='#FFFFFF', align=PP_ALIGN.CENTER)
        add_text(s5, task, cx+Inches(0.54), cy+Inches(0.66)+j*Inches(1.1),
                 col_w-Inches(0.65), Inches(0.72), size=11.5, color=C['dark'])

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — AI 月報一條龍
# ─────────────────────────────────────────────────────────────────────────────
s6 = prs.slides.add_slide(blank_layout)
set_slide_bg(s6, '#0F172A')

# header
add_rect(s6, Inches(0), Inches(0), W, Inches(1.35), fill='#064E3B')
add_text(s6, 'AI 驅動月報一條龍', Inches(0.45), Inches(0.12), Inches(9), Inches(0.72),
         size=26, bold=True, color='#4ADE80')
add_text(s6, '從原始資料到簡報，全程 AI 工具協助完成', Inches(0.45), Inches(0.82), Inches(9), Inches(0.4),
         size=13, color='#86EFAC')
add_rect(s6, Inches(0), Inches(1.35), W, Inches(0.04), fill='#22C55E')

# AI badge top right
add_rect(s6, Inches(10.5), Inches(0.2), Inches(2.6), Inches(0.55), fill='#166534', line='#4ADE80', line_w=Pt(1.5))
add_text(s6, '🤖 AI-First 工作流程', Inches(10.5), Inches(0.2), Inches(2.6), Inches(0.55),
         size=11, bold=True, color='#4ADE80', align=PP_ALIGN.CENTER)

# flow steps
steps = [
    ('📊', '原始資料', 'Excel / PDF\n工作報告收集', '#1E3A5F', '#3B82F6'),
    ('🤖', 'AI 分析', 'Claude AI 提取\n結構化資料', '#064E3B', '#22C55E'),
    ('📈', '互動儀表板', '自動生成月報\nWeb Dashboard', '#3B0764', '#A855F7'),
    ('📋', '部門月報', 'KPI 追蹤\n案件明細一覽', '#451A03', '#F59E0B'),
    ('🎯', '本份簡報', 'AI 自動排版\n簡報一條龍', '#4C0519', '#EF4444'),
]

sw = Inches(2.3)
for i,(icon,title,desc,bg,col) in enumerate(steps):
    sx = Inches(0.35) + i*(sw+Inches(0.25))
    sy = Inches(1.65)
    add_rect(s6, sx, sy, sw, Inches(2.8), fill=bg, line=col, line_w=Pt(1.5))
    add_text(s6, icon, sx+Inches(0.1), sy+Inches(0.2), sw-Inches(0.2), Inches(0.65),
             size=30, align=PP_ALIGN.CENTER)
    add_text(s6, title, sx+Inches(0.1), sy+Inches(0.85), sw-Inches(0.2), Inches(0.45),
             size=16, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s6, desc, sx+Inches(0.15), sy+Inches(1.3), sw-Inches(0.3), Inches(0.85),
             size=11, color='#CBD5E1', align=PP_ALIGN.CENTER)
    # arrow (not last)
    if i < 4:
        add_text(s6, '▶', sx+sw+Inches(0.04), sy+Inches(1.05), Inches(0.22), Inches(0.45),
                 size=14, color='#475569', align=PP_ALIGN.CENTER)

# benefits row
benefits = [
    ('⏱️ 效率提升', '報告製作時間\n縮短 80%'),
    ('🎯 零排版', '自動生成，無需\n人工美化調整'),
    ('🔄 即時更新', '資料異動後\n一鍵重新生成'),
    ('📊 數據驅動', '視覺化圖表\n自動彙整分析'),
]
bw2 = Inches(2.9)
for i,(title,desc) in enumerate(benefits):
    bx = Inches(0.6) + i*(bw2+Inches(0.35))
    add_rect(s6, bx, Inches(4.85), bw2, Inches(1.0), fill='#1E293B', line='#334155', line_w=Pt(1))
    add_text(s6, title, bx+Inches(0.12), Inches(4.9), bw2-Inches(0.24), Inches(0.38),
             size=13, bold=True, color='#4ADE80')
    add_text(s6, desc, bx+Inches(0.12), Inches(5.25), bw2-Inches(0.24), Inches(0.5),
             size=10.5, color='#94A3B8')

add_text(s6, '本份月報與簡報均為 AI 工具（Claude）協助完成，展示部門 AI 工具應用實力',
         Inches(0.35), Inches(6.15), Inches(12.6), Inches(0.45),
         size=11, italic=True, color='#475569', align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — 年度 KPI 進度
# ─────────────────────────────────────────────────────────────────────────────
s7 = prs.slides.add_slide(blank_layout)
set_slide_bg(s7, '#F8FAFC')
slide_header(s7, '2026 年度 KPI 進度追蹤', 'Q1–Q2 重點指標現況')

kpi_items = [
    ('月均技術支援案件達標', 'Denny 117件 ✓  William 43件 ✓  Lucas 51件 ✓  Jimmy 72件 ✓', '#22C55E', '進行中'),
    ('收費服務成交目標', 'William NT$75,000 實收 ✓  其他成員積極推進中', '#3B82F6', '進行中'),
    ('對外 Workshop 年度目標', 'Jimmy 雲科大 Firefly AI 講座 ✓  年度目標 ≥8 次持續推進', '#A855F7', '進行中'),
    ('知識庫文件撰寫', 'Denny / William 季度目標推進中  Jimmy KB 文章定期更新', '#F59E0B', '進行中'),
    ('原廠認證更新', '各成員依代理商要求維持有效認證，年度前完成 renew', '#0891B2', '待執行'),
    ('AI 工具導入', 'Denny / William 已應用於案件回覆與文件生成  Jimmy Firefly 整合推廣', '#22C55E', '達成'),
    ('部門月報呈報', '每月按時彙整並透過 AI 互動儀表板與簡報呈報 ✓', '#22C55E', '達成'),
]

status_style = {
    '達成':  ('#D1FAE5', '#065F46'),
    '進行中':('#DBEAFE', '#1D4ED8'),
    '待執行':('#FEF3C7', '#92400E'),
}

for i,(title,detail,col,status) in enumerate(kpi_items):
    row_y = Inches(1.65) + i*Inches(0.72)
    # left color bar
    add_rect(s7, Inches(0.2), row_y, Inches(0.08), Inches(0.58), fill=col)
    # row bg
    add_rect(s7, Inches(0.32), row_y, Inches(12.6), Inches(0.58),
             fill='#FFFFFF' if i%2==0 else '#F8FAFC', line='#E2E8F0', line_w=Pt(0.5))
    add_text(s7, title, Inches(0.48), row_y+Inches(0.06), Inches(3.5), Inches(0.42),
             size=12, bold=True, color=C['dark'])
    add_text(s7, detail, Inches(4.1), row_y+Inches(0.08), Inches(7.7), Inches(0.42),
             size=10.5, color='#475569')
    # status badge
    sb, sc = status_style.get(status, ('#F1F5F9','#475569'))
    add_rect(s7, Inches(12.0), row_y+Inches(0.1), Inches(0.9), Inches(0.32), fill=sb)
    add_text(s7, status, Inches(12.0), row_y+Inches(0.1), Inches(0.9), Inches(0.32),
             size=9.5, bold=True, color=sc, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — 結語 & Next Action
# ─────────────────────────────────────────────────────────────────────────────
s8 = prs.slides.add_slide(blank_layout)
set_slide_bg(s8, '#0F172A')

add_rect(s8, Inches(0), Inches(0), Inches(0.18), H, fill='#0891B2')

add_text(s8, '結語 & 下月目標', Inches(0.55), Inches(0.8), Inches(10), Inches(0.65),
         size=28, bold=True, color='#FFFFFF')
add_rect(s8, Inches(0.55), Inches(1.48), Inches(5), Inches(0.04), fill='#38BDF8')

# summary points
summary_pts = [
    ('📈', '業務持續成長', '本月 379 件案件，月均虛擬估值 NT$216,500，William 實收 NT$75,000，整體績效穩健。'),
    ('🚀', '5月重點衝刺', '彰銀 Jamf 憑證更新、奇美醫院進階培訓、Sentinel POC 驗收、Firefly AI 課程上線。'),
    ('🤖', 'AI 工具持續落地', '月報一條龍 + 簡報自動化已成功驗證，下月規劃擴展至客戶提案文件自動化。'),
    ('🏆', 'KPI 進度健康', '月均案件、AI 工具導入、月報呈報三項已達標；收費服務持續累積中。'),
]

for i,(icon,title,desc) in enumerate(summary_pts):
    ry = Inches(1.75) + i*Inches(1.25)
    add_rect(s8, Inches(0.55), ry, Inches(11.5), Inches(1.05), fill='#1E293B', line='#334155', line_w=Pt(1))
    add_text(s8, icon, Inches(0.65), ry+Inches(0.05), Inches(0.65), Inches(0.65),
             size=26, align=PP_ALIGN.CENTER)
    add_text(s8, title, Inches(1.35), ry+Inches(0.08), Inches(4), Inches(0.4),
             size=14, bold=True, color='#38BDF8')
    add_text(s8, desc, Inches(1.35), ry+Inches(0.5), Inches(10.5), Inches(0.5),
             size=12, color='#94A3B8')

# footer
add_rect(s8, Inches(0), Inches(7.08), W, Inches(0.42), fill='#064E3B')
add_text(s8, '🤖 本月報 / 儀表板 / 簡報均由 AI 工具協助完成  ‧  KS010S 部門  ‧  2026年5月',
         Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.38),
         size=10.5, color='#4ADE80', align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────────────────────────────────────
out = r'D:\excel_dashboard07 月報\KS010S_2026年4月月報簡報.pptx'
prs.save(out)
print(f'Saved: {out}')
